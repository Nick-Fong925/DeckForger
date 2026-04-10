import csv
import os
import re
import sqlite3
import tempfile
import uuid
import zipfile
from html import unescape
from typing import Any

MAX_CARDS_PER_DECK = 500

_CLOZE_RE = re.compile(r'\{\{c\d+::([^:}]+)(?:::[^}]*)?\}\}')
_HEADER_KEYWORDS = {'front', 'back', 'question', 'answer', 'term', 'definition'}


def _strip_html(text: str) -> str:
    text = re.sub(r'\[sound:[^\]]+\]', '', text)  # Anki [sound:...] tags
    text = re.sub(r'<[^>]+>', '', text)            # HTML tags
    return unescape(text).strip()                  # HTML entities (&amp; etc.)


def _clean_anki_html(text: str) -> str:
    """Remove Anki-specific non-HTML tokens but preserve HTML for frontend rendering."""
    return re.sub(r'\[sound:[^\]]+\]', '', text).strip()


def _cloze_to_qa(text: str) -> tuple[str, str]:
    answers = _CLOZE_RE.findall(text)
    front = _CLOZE_RE.sub('_____', text)
    return front, ', '.join(answers)


def _is_header_row(row: list[str]) -> bool:
    return len(row) >= 2 and row[0].strip().lower() in _HEADER_KEYWORDS


def pdf_to_text(file_path: str) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    pages = [page.get_text() for page in doc]
    doc.close()
    return '\n\n'.join(pages)


def pptx_to_text(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    slides: list[str] = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if shape.has_text_frame]
        slides.append('\n'.join(texts))
    return '\n\n'.join(slides)


def csv_to_cards(file_path: str, delimiter: str = ',') -> list[dict[str, Any]]:
    cards: list[dict[str, Any]] = []
    # utf-8-sig handles Windows BOM; falls back gracefully for standard UTF-8
    with open(file_path, newline='', encoding='utf-8-sig') as f:
        reader = csv.reader(f, delimiter=delimiter)
        first_row = next(reader, None)
        if first_row is None:
            return cards
        # Include first row as a card only if it doesn't look like a header
        rows_to_process = iter([first_row]) if not _is_header_row(first_row) else iter([])
        for row in [*rows_to_process, *reader]:
            if len(cards) >= MAX_CARDS_PER_DECK:
                break
            if len(row) < 2:
                continue
            front = row[0].strip()
            back = row[1].strip()
            if not front or not back:
                continue
            cards.append({'id': str(uuid.uuid4()), 'front': front, 'back': back})
    return cards


def tsv_to_cards(file_path: str) -> list[dict[str, Any]]:
    return csv_to_cards(file_path, delimiter='\t')


def _open_anki_db(tmp_dir: str) -> sqlite3.Connection:
    # Prefer newer formats first; fall back to legacy collection.anki2.
    # collection.anki21b: zstd-compressed SQLite (Anki 2.1.50+)
    anki21b = os.path.join(tmp_dir, 'collection.anki21b')
    if os.path.exists(anki21b):
        import zstandard
        with open(anki21b, 'rb') as f:
            compressed = f.read()
        data = zstandard.ZstdDecompressor().decompress(compressed, max_output_size=500 * 1024 * 1024)
        out_path = os.path.join(tmp_dir, 'collection.sqlite')
        with open(out_path, 'wb') as f:
            f.write(data)
        return sqlite3.connect(out_path)

    # collection.anki21: plain SQLite (Anki 2.1.x before 2.1.50)
    anki21 = os.path.join(tmp_dir, 'collection.anki21')
    if os.path.exists(anki21):
        return sqlite3.connect(anki21)

    # collection.anki2: legacy SQLite (Anki 2.0 / stub in newer exports)
    return sqlite3.connect(os.path.join(tmp_dir, 'collection.anki2'))


def apkg_to_cards(file_path: str) -> list[dict[str, Any]]:
    # .apkg is a zip containing a SQLite database with notes.
    # The `flds` column uses \x1f (unit separator) to delimit fields.
    # Field content is stored as HTML; strip tags before returning.
    cards: list[dict[str, Any]] = []
    with tempfile.TemporaryDirectory() as tmp_dir:
        with zipfile.ZipFile(file_path, 'r') as zf:
            for member in zf.namelist():
                member_path = os.path.realpath(os.path.join(tmp_dir, member))
                if not member_path.startswith(os.path.realpath(tmp_dir) + os.sep):
                    raise ValueError(f'Zip slip attempt detected in member: {member}')
            zf.extractall(tmp_dir)
        conn = _open_anki_db(tmp_dir)
        try:
            cursor = conn.execute('SELECT flds FROM notes')
            for (flds,) in cursor:
                if len(cards) >= MAX_CARDS_PER_DECK:
                    break
                parts = flds.split('\x1f')
                if len(parts) < 2:
                    continue
                # Preserve HTML so the frontend can render formatting, line breaks, etc.
                # Only strip Anki-specific [sound:...] tokens — DOMPurify sanitizes on the client.
                raw_front = _clean_anki_html(parts[0])
                if _CLOZE_RE.search(raw_front):
                    front, back = _cloze_to_qa(raw_front)
                else:
                    front = raw_front
                    back = _clean_anki_html(parts[1])
                # Use tag-stripped version only for the emptiness check
                if not _strip_html(front) or not _strip_html(back):
                    continue
                cards.append({'id': str(uuid.uuid4()), 'front': front, 'back': back})
        finally:
            conn.close()
    return cards
